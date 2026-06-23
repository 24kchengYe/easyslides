# -*- coding: utf-8 -*-
"""
Reusable python-pptx helpers for the AI4US talk comprehensibility pass.

- add_gloss_box(): a small, visually-distinct "plain-language" caption box
  matching the deck's existing card style (Segoe UI; soft tint; rounded).
- add_key_slide(): a single plain-language glossary slide ("How to read this talk").
- GLOSS: canonical term -> short plain-language definition for a NON-AI
  urban-planning professor. Numbers/claims must stay faithful to the paper.

All edits are IN PLACE. Never re-export via easyslides.
Deck style: title Segoe UI 25.5 bold #1A1A1A; subtitle 12.75 #666666;
pagenum 9 grey bottom-right; blue accent bar #0072B2; white bg.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# palette (matches deck + figstyle)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x66, 0x66, 0x66)
BLUE   = RGBColor(0x00, 0x72, 0xB2)
GREEN  = RGBColor(0x00, 0x9E, 0x73)
ORANGE = RGBColor(0xD5, 0x5E, 0x00)
RED    = RGBColor(0xC4, 0x4E, 0x52)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TINT   = RGBColor(0xEA, 0xF3, 0xF8)   # light blue tint (deck's pill colour)
TINT_BORDER = RGBColor(0xBF, 0xD7, 0xE6)
DARK   = RGBColor(0x13, 0x2A, 0x3E)   # deck's dark banner
DARK_LABEL = RGBColor(0x9F, 0xB4, 0xC4)

FONT = "Segoe UI"


def _set_no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def add_gloss_box(slide, left, top, width, height, label, body,
                  accent=BLUE, fill=TINT, border=TINT_BORDER,
                  label_size=10.5, body_size=10.5, body_color=INK):
    """A small rounded annotation box: bold accent label + plain-language body.
    Placed where the audited slide has free space. Sizes in inches.
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border; box.line.width = Pt(0.75)
    _set_no_shadow(box)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.name = FONT; r.font.size = Pt(label_size)
        r.font.bold = True; r.font.color.rgb = accent
    r2 = p.add_run(); r2.text = body
    r2.font.name = FONT; r2.font.size = Pt(body_size)
    r2.font.bold = False; r2.font.color.rgb = body_color
    return box


def add_plain_caption(slide, left, top, width, height, body,
                      size=10.5, color=GREY, italic=True, align=PP_ALIGN.LEFT):
    """A borderless small caption (no box) for under-figure plain-language notes."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = body
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def _title_block(slide, title, subtitle, SW, SH):
    """Clone the deck's title block: white bg + blue accent bar + title + subtitle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    _set_no_shadow(bg)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.625), Inches(0.5),
                                 Inches(0.062), Inches(0.458))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    _set_no_shadow(bar)
    tb = slide.shapes.add_textbox(Inches(0.819), Inches(0.449), Inches(11.6), Inches(0.567))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.name = FONT; r.font.size = Pt(25.5); r.font.bold = True; r.font.color.rgb = INK
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.836), Inches(0.945), Inches(11.9), Inches(0.3))
        stf = sb.text_frame; stf.word_wrap = True
        sp = stf.paragraphs[0]; sp.alignment = PP_ALIGN.LEFT
        sr = sp.add_run(); sr.text = subtitle
        sr.font.name = FONT; sr.font.size = Pt(12.75); sr.font.bold = False; sr.font.color.rgb = GREY


def add_pagenum(slide, num, color=GREY):
    tb = slide.shapes.add_textbox(Inches(12.329), Inches(7.185), Inches(0.3), Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{num:02d}"
    r.font.name = FONT; r.font.size = Pt(9); r.font.color.rgb = color


# canonical term -> (short label, plain-language definition)
# tight, audience = urban-planning professor; faithful to the paper.
GLOSS = {
    "probe":      "we fit a simple ruler to the model's internal numbers to test whether a fact is encoded inside",
    "intervene":  "we directly nudge those internal numbers and check whether the model's answer changes (a true cause-test)",
    "activation": "the internal numbers the model computes as it reads a prompt (its working 'thoughts')",
    "hidden_state": "the internal numbers (activations) the model holds at one layer; written h",
    "layer":      "a model is a stack of processing steps called layers; early = near input, deep = near the answer",
    "sweep":      "we repeat the test at every layer, one by one, to see where the effect peaks",
    "causal_layer": "the layer where nudging the encoded fact actually changes the output",
    "memory_layer": "the layer where the fact is most readable; storing a fact is not the same as using it",
    "layer_localized": "the cause-effect shows up in just one or two layers, not spread across the model",
    "R2":         "goodness of fit, 0 to 1; higher means the internal numbers track the real value more tightly",
    "kappa":      "Cohen's Kappa, human-vs-AI agreement; 0 = chance, 1 = perfect (0.16-0.30 = weak-to-fair)",
    "moran":      "Moran's I, a spatial-clustering score; positive means nearby places get similar ratings",
    "tobler":     "Tobler's first law of geography: near things are more alike than distant things",
    "VL":         "Vision-Language model, reads images as well as text (here: open-weight Qwen2.5-VL)",
    "open_closed": "closed models (GPT-4o) we can only prompt; open models (Qwen) we can read and edit inside",
    "w_pop":      "the direction inside the model that encodes population, found by the probe",
    "delta":      "a small push we add along that direction to test cause and effect",
    "logodds":    "log-odds: + values make the high-density answer more likely; the curve's slope is the effect size",
    "relative_depth": "layer position as a fraction: 0 = first layer, 1 = last layer (0.5 = middle)",
    "FFN":        "feed-forward block, the part of the model that stores facts like a look-up table",
    "MAE":        "mean absolute error, average size of the miss; lower is better",
    "JSD":        "a 0-to-1 distance between two distributions; 0 = identical, larger = more different",
    "CLIP":       "an image-similarity score; higher = more varied street scenes",
    "OLS":        "ordinary least-squares regression; beta is the slope = strength and direction of each effect",
    "cross_val":  "tested on held-out cities the probe never saw, so the fit is not memorised",
    "sign_agree": "how many effects point the same way (+/-) in both models",
}


def add_key_slide(prs, after_index, pagenum, pagenum_color=GREY, level="full"):
    """Insert a single plain-language 'how to read this talk' glossary slide
    right AFTER `after_index` (0-based). Renumber handled by caller.
    level: 'full' (60min) or 'medium' (30min) selects how many terms to show.
    Returns the new slide.
    """
    SW, SH = prs.slide_width, prs.slide_height
    blank = next((l for l in prs.slide_layouts if l.name == 'Blank'), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    _title_block(slide, "How to read this talk: a plain-language key",
                 "A quick glossary of the few technical terms used in the figures, in everyday language",
                 SW, SH)

    # term cards in 2 columns. Each: bold term + plain definition.
    if level == "medium":
        items = [
            ("Probe", GLOSS["probe"]),
            ("Intervene (cause-test)", GLOSS["intervene"]),
            ("Layer / layer sweep", "a model is a stack of steps (layers); we test every layer to see where an effect peaks"),
            ("Memory layer vs causal layer", "where a fact is readable vs where changing it changes the answer; they need not coincide"),
            ("Open vs closed model", GLOSS["open_closed"]),
            ("VL (vision-language)", GLOSS["VL"]),
            ("R squared", GLOSS["R2"]),
            ("Kappa (perception)", GLOSS["kappa"]),
            ("Moran's I / Tobler", "spatial-clustering score; positive = nearby places rated alike (Tobler's law of geography)"),
        ]
    else:  # full
        items = [
            ("Probe", GLOSS["probe"]),
            ("Intervene (cause-test)", GLOSS["intervene"]),
            ("Hidden state h", GLOSS["hidden_state"]),
            ("Layer / sweep", "the model is a stack of steps (layers); a sweep tests every layer one by one"),
            ("Relative depth", GLOSS["relative_depth"]),
            ("Memory vs causal layer", "where a fact is readable vs where changing it changes the answer"),
            ("Open vs closed model", GLOSS["open_closed"]),
            ("VL (vision-language)", GLOSS["VL"]),
            ("R squared", GLOSS["R2"]),
            ("Kappa", GLOSS["kappa"]),
            ("Moran's I / Tobler", "spatial-clustering score; positive = nearby places rated alike"),
            ("Effect size (log-odds)", GLOSS["logodds"]),
        ]

    n = len(items)
    rows = (n + 1) // 2
    col_w = 5.75
    x0, x1 = 0.82, 6.74
    y_top = 1.55
    row_h = (6.95 - y_top) / rows
    for k, (term, defn) in enumerate(items):
        col = k // rows
        row = k % rows
        x = x0 if col == 0 else x1
        y = y_top + row * row_h
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y + 0.04), Inches(col_w), Inches(row_h - 0.12))
        box.fill.solid(); box.fill.fore_color.rgb = TINT
        box.line.color.rgb = TINT_BORDER; box.line.width = Pt(0.75)
        _set_no_shadow(box)
        tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = term + " — "
        r.font.name = FONT; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = BLUE
        r2 = p.add_run(); r2.text = defn
        r2.font.name = FONT; r2.font.size = Pt(11.5); r2.font.color.rgb = INK

    add_pagenum(slide, pagenum, pagenum_color)

    # move into position after `after_index`
    lst = prs.slides._sldIdLst
    moved = list(lst)[-1]
    lst.remove(moved)
    lst.insert(after_index + 1, moved)
    return slide


def add_divider_slide(prs, after_index, label, headline_top, headline_accent,
                      sub, accent=BLUE, pagenum=None, pagenum_color=None,
                      footer=None):
    """Insert a dark section-divider slide (132A3E bg) AFTER `after_index` (0-based),
    matching the deck's PART-divider recipe. Returns the new slide.

    label:           small pale-blue caps label (e.g. 'LAYER 1 . SURFACE')
    headline_top:    white headline line(s); list of strings (one per line) or str
    headline_accent: optional final headline line in accent colour (str or None)
    sub:             italic pale tagline under the headline
    accent:          accent-bar + emphasis colour (BLUE/ORANGE/RED)
    footer:          optional small pale footer line near the bottom
    """
    SW, SH = prs.slide_width, prs.slide_height
    blank = next((l for l in prs.slide_layouts if l.name == 'Blank'), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    PALE = RGBColor(0x9F, 0xB4, 0xC4)
    TAG = RGBColor(0xCB, 0xD8, 0xE2)

    # dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
    _set_no_shadow(bg)
    # top accent bar (full width, thin)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.0625))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    _set_no_shadow(bar)

    # label (centered)
    lb = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.33), Inches(0.35))
    lb.text_frame.word_wrap = True
    lp = lb.text_frame.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run(); lr.text = label
    lr.font.name = FONT; lr.font.size = Pt(13.5); lr.font.bold = True; lr.font.color.rgb = PALE

    # headline lines (white) + optional accent line, centered, large
    lines = headline_top if isinstance(headline_top, (list, tuple)) else [headline_top]
    y = 2.45
    for ln in lines:
        hb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.73), Inches(0.85))
        hb.text_frame.word_wrap = True
        hp = hb.text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = ln
        hr.font.name = FONT; hr.font.size = Pt(40); hr.font.bold = True; hr.font.color.rgb = WHITE
        y += 0.78
    if headline_accent:
        hb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.73), Inches(0.85))
        hb.text_frame.word_wrap = True
        hp = hb.text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = headline_accent
        hr.font.name = FONT; hr.font.size = Pt(40); hr.font.bold = True; hr.font.color.rgb = accent
        y += 0.78

    # tagline
    if sub:
        sb = slide.shapes.add_textbox(Inches(1.0), Inches(max(y + 0.15, 5.0)), Inches(11.33), Inches(0.5))
        sb.text_frame.word_wrap = True
        spp = sb.text_frame.paragraphs[0]; spp.alignment = PP_ALIGN.CENTER
        sr = spp.add_run(); sr.text = sub
        sr.font.name = FONT; sr.font.size = Pt(16); sr.font.italic = True; sr.font.color.rgb = TAG

    if footer:
        fb = slide.shapes.add_textbox(Inches(1.0), Inches(6.45), Inches(11.33), Inches(0.4))
        fb.text_frame.word_wrap = True
        fp = fb.text_frame.paragraphs[0]; fp.alignment = PP_ALIGN.CENTER
        fr = fp.add_run(); fr.text = footer
        fr.font.name = FONT; fr.font.size = Pt(13.5); fr.font.color.rgb = PALE

    if pagenum is not None:
        add_pagenum(slide, pagenum, pagenum_color or GREY)

    lst = prs.slides._sldIdLst
    moved = list(lst)[-1]
    lst.remove(moved)
    lst.insert(after_index + 1, moved)
    return slide


def add_titled_blank_slide(prs, after_index, title, subtitle, pagenum=None,
                           pagenum_color=None):
    """Insert a blank white slide with the standard title block AFTER `after_index`.
    Returns the new slide (caller can then inject an SVG diagram or add shapes).
    """
    SW, SH = prs.slide_width, prs.slide_height
    blank = next((l for l in prs.slide_layouts if l.name == 'Blank'), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    _title_block(slide, title, subtitle, SW, SH)
    if pagenum is not None:
        add_pagenum(slide, pagenum, pagenum_color or GREY)
    lst = prs.slides._sldIdLst
    moved = list(lst)[-1]
    lst.remove(moved)
    lst.insert(after_index + 1, moved)
    return slide


def remove_slide(prs, index):
    """Remove the slide at 0-based `index` properly: drop the sldId entry AND the
    presentation->slide relationship, so the part is not resurrected on save."""
    lst = prs.slides._sldIdLst
    sldIds = list(lst)
    sldId = sldIds[index]
    rId = sldId.get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    # drop the relationship from the presentation part
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    lst.remove(sldId)


def renumber_pages(prs):
    """Renumber bottom-right 2-digit page numbers to match slide order."""
    for i, s in enumerate(prs.slides, 1):
        if i == 1:
            continue
        for sh in s.shapes:
            if sh.has_text_frame and sh.left is not None:
                L = Emu(sh.left).inches; T = Emu(sh.top).inches
                t = sh.text_frame.text.strip()
                if T > 6.9 and L > 11.5 and t.isdigit() and len(t) <= 4:
                    para = sh.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = f"{i:02d}"
                        for extra in para.runs[1:]:
                            extra.text = ""
                    break
